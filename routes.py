# routes.py
import os
import json
import matplotlib
matplotlib.use('Agg')  # Verwenden Sie 'Agg' Backend für nicht-interaktive Umgebungen
from flask import Blueprint, render_template, redirect, url_for, request, flash
from forms import ProjectForm, create_edit_questions_form
from models import Project
from extensions import db
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
from io import BytesIO
from flask import send_file
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY
from reportlab.lib import colors
import numpy as np
import logging
from flask import send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


routes_bp = Blueprint('routes', __name__)

# Logging konfigurieren
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

meddpicc_elements = [
    {'title': 'Metrics', 'text': 'Key metrics to measure the success of the deal.', 'route': 'metrics'},
    {'title': 'Economic Buyer', 'text': 'Identify the key decision-maker who controls the budget.', 'route': 'economic_buyer'},
    {'title': 'Decision Criteria', 'text': 'Understand the criteria that will be used to make the decision.', 'route': 'decision_criteria'},
    {'title': 'Decision Process', 'text': 'Know the decision-making process within the company.', 'route': 'decision_process'},
    {'title': 'Paper Process', 'text': 'Understand the contract closing process.', 'route': 'paper_process'},
    {'title': 'Identify Pain', 'text': 'Identify the customer\'s pain points and challenges.', 'route': 'identify_pain'},
    {'title': 'Champion', 'text': 'Find an internal supporter within the customer\'s company.', 'route': 'champion'},
    {'title': 'Competition', 'text': 'Analyze the competitive landscape.', 'route': 'competition'}
]

# Intro route
@routes_bp.route('/')
def intro():
    return render_template('intro.html', meddpicc_elements=meddpicc_elements)

# Metrics.html route
@routes_bp.route('/metrics')
def metrics():
    return render_template('metrics.html')

# economic_buyer.html route
@routes_bp.route('/economic-buyer')
def economic_buyer():
    return render_template('economic_buyer.html')

# decision_criteria.html route
@routes_bp.route('/decision_criteria')
def decision_criteria():
    return render_template('decision_criteria.html')

# decision_process.html route
@routes_bp.route('/decision_process')
def decision_process():
    return render_template('decision_process.html')

# paper_process.html route
@routes_bp.route('/paper_process')
def paper_process():
    return render_template('paper_process.html')

# identify_pain.html route
@routes_bp.route('/identify_pain')
def identify_pain():
    return render_template('identify_pain.html')

# champion.html route
@routes_bp.route('/champion')
def champion():
    return render_template('champion.html')

# competition.html route
@routes_bp.route('/competition')
def competition():
    return render_template('competition.html')

# Index route
@routes_bp.route('/index')
def index():
    projects = Project.query.all()
    return render_template('index.html', projects=projects)

# Route to add a new project
@routes_bp.route('/new_project', methods=['GET', 'POST'])
def new_project():
    form = ProjectForm()
    if form.validate_on_submit():
        project = Project()
        form.populate_obj(project)
        
        # Logging der Projektfelder
        for field in form:
            if field.name != 'csrf_token' and field.name != 'submit':
                logger.info(f"{field.name}: {field.data}")
        
        # Berechnen, ob Fragen beantwortet wurden (Ja/Nein)
        for element in ['metrics', 'economic_buyer', 'decision_criteria', 'decision_process', 'paper_process', 'implications_of_pain', 'champion', 'competition']:
            answered = False
            for q_num in [1, 2, 3]:
                field_name = f"{element}_question{q_num}"
                if getattr(project, field_name):
                    answered = True
                    break
            setattr(project, f"{element}_question_answered", answered)  # BooleanField im Modell

        db.session.add(project)
        db.session.commit()
        flash('Neues Projekt erfolgreich hinzugefügt!', 'success')
        return redirect(url_for('routes.index'))
    return render_template('new_project.html', form=form)

# Route to edit an existing project
@routes_bp.route('/edit_project/<int:project_id>', methods=['GET', 'POST'])
def edit_project(project_id):
    project = Project.query.get_or_404(project_id)
    form = ProjectForm(obj=project)
    if form.validate_on_submit():
        form.populate_obj(project)
        
        # Logging der Projektfelder
        for field in form:
            if field.name not in ['csrf_token', 'submit']:
                logger.info(f"{field.name}: {field.data}")
        
        # Berechnen, ob Fragen beantwortet wurden (Ja/Nein)
        for element in ['metrics', 'economic_buyer', 'decision_criteria', 'decision_process', 'paper_process', 'implications_of_pain', 'champion', 'competition']:
            answered = False
            for q_num in [1, 2, 3]:
                field_name = f"{element}_question{q_num}"
                value = getattr(project, field_name)
                logger.info(f"Überprüfe {field_name}: {value}")
                if value and value.strip():  # Überprüfen, ob das Feld nicht leer ist
                    answered = True
                    break
            setattr(project, f"{element}_question_answered", answered)
            logger.info(f"Setze {element}_question_answered auf {answered}")

        db.session.add(project)
        db.session.commit()
        flash('Projekt erfolgreich aktualisiert!', 'success')
        return redirect(url_for('routes.index'))
    return render_template('edit_project.html', form=form, project=project)

# Route to delete a project
@routes_bp.route('/delete_project/<int:project_id>', methods=['POST'])
def delete_project(project_id):
    project = Project.query.get_or_404(project_id)
    db.session.delete(project)
    db.session.commit()
    flash('Project deleted successfully!', 'success')
    return redirect(url_for('routes.index'))

# Route for the radar chart
@routes_bp.route('/radar/<int:project_id>')
def radar(project_id):
    project = Project.query.get_or_404(project_id)
    
    # Daten für das Radar-Diagramm
    data = {
        'Metrics': project.metrics,
        'Economic Buyer': project.economic_buyer,
        'Decision Criteria': project.decision_criteria,
        'Decision Process': project.decision_process,
        'Paper Process': project.paper_process,
        'Implications of Pain': project.implications_of_pain,
        'Champion': project.champion,
        'Competition': project.competition
    }
    
    categories = list(data.keys())
    values = list(data.values())
    values += values[:1]  # Radar-Diagramm schließen
    
    angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
    angles += angles[:1]  # Schließen
    
    fig, ax = plt.subplots(figsize=(6,6), subplot_kw=dict(polar=True))
    
    ax.plot(angles, values, color='blue', linewidth=2, linestyle='solid')
    ax.fill(angles, values, color='blue', alpha=0.25)
    
    ax.set_yticklabels([])
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories)
    
    plt.title(f'Radar Chart for {project.name}')
    
    radar_path = os.path.join('static', 'images', f'radar_{project_id}.png')
    plt.savefig(radar_path)
    plt.close()
    
    return render_template('radar.html', project=project)


# Route for the scorecard
@routes_bp.route('/scorecard/<int:project_id>')
def scorecard(project_id):
    project = Project.query.get_or_404(project_id)
    
    # Daten für die Scorecard
    data = {
        'Metrics': [project.metrics],
        'Economic Buyer': [project.economic_buyer],
        'Decision Criteria': [project.decision_criteria],
        'Decision Process': [project.decision_process],
        'Paper Process': [project.paper_process],
        'Implications of Pain': [project.implications_of_pain],
        'Champion': [project.champion],
        'Competition': [project.competition]
    }
    
    df = pd.DataFrame(data)
    sum_scores = df.sum().to_frame(name='Sum Score').T

    # Sicherstellen, dass das Verzeichnis existiert
    images_dir = os.path.join('static', 'images')
    os.makedirs(images_dir, exist_ok=True)

    # Scorecard rendering with seaborn
    plt.figure(figsize=(12, 2))
    sns.heatmap(sum_scores, annot=True, cmap='RdYlGn', cbar=False)
    plt.title(f'Sum Scores Scorecard for {project.name}')
    scorecard_path = os.path.join(images_dir, f'scorecard_{project_id}.png')
    plt.savefig(scorecard_path)
    plt.close()

    return render_template('scorecard.html', project=project)

# Route to manage and display questions from questions.json
@routes_bp.route('/manage_questions', methods=['GET'])
def manage_questions():
    questions_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static', 'questions.json')
    try:
        with open(questions_path, 'r', encoding='utf-8') as f:
            questions = json.load(f)
    except FileNotFoundError:
        flash("Die questions.json Datei wurde nicht gefunden.", "danger")
        questions = {}
    except json.JSONDecodeError:
        flash("Fehler beim Parsen der questions.json Datei.", "danger")
        questions = {}
    
    return render_template('questions.html', questions=questions)

# Route to edit questions
@routes_bp.route('/edit_questions', methods=['GET', 'POST'])
def edit_questions():
    questions_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static', 'questions.json')
    try:
        with open(questions_path, 'r', encoding='utf-8') as f:
            questions = json.load(f)
    except FileNotFoundError:
        flash("Die questions.json Datei wurde nicht gefunden.", "danger")
        questions = {}
    except json.JSONDecodeError:
        flash("Fehler beim Parsen der questions.json Datei.", "danger")
        questions = {}
    
    # Dynamisch die Formularklasse erstellen und instanziieren
    EditQuestionsForm = create_edit_questions_form()
    form = EditQuestionsForm()
    
    if request.method == 'GET':
        # Initialisiere die Formulardaten mit den aktuellen Fragen
        for element, element_questions in questions.items():
            for q_key, q_text in element_questions.items():
                field_name = f"{element}_{q_key}"
                if hasattr(form, field_name):
                    getattr(form, field_name).data = q_text
    
    if form.validate_on_submit():
        # Aktualisieren der Fragen basierend auf den Formulardaten
        for element, element_questions in questions.items():
            for q_key in element_questions.keys():
                field_name = f"{element}_{q_key}"
                new_text = getattr(form, field_name).data
                if new_text and new_text.strip():
                    questions[element][q_key] = new_text.strip()
                else:
                    flash(f"Frage {q_key} in {element.replace('_', ' ').capitalize()} darf nicht leer sein.", "danger")
                    return redirect(url_for('routes.edit_questions'))
        
        # Speichern der aktualisierten Fragen zurück in questions.json
        try:
            with open(questions_path, 'w', encoding='utf-8') as f:
                json.dump(questions, f, ensure_ascii=False, indent=4)
            flash("Fragen erfolgreich aktualisiert!", "success")
            return redirect(url_for('routes.manage_questions'))
        except Exception as e:
            logger.error(f"Fehler beim Schreiben in questions.json: {e}")
            flash("Fehler beim Speichern der Fragen.", "danger")
            return redirect(url_for('routes.edit_questions'))
    
    return render_template('edit_questions.html', form=form, questions=questions)

# Neue Funktion für PDF-Export
def create_scorecard_table(project):
    elements = ['metrics', 'economic_buyer', 'decision_criteria', 'decision_process', 'paper_process',
                'implications_of_pain', 'champion', 'competition']
    
    data = [['Element', 'Fragen beantwortet', 'Score (1-10)', 'Kommentare']]
    
    for element in elements:
        element_name = element.replace('_', ' ').capitalize()
        questions_answered = 'Ja' if getattr(project, f'{element}_question_answered', False) else 'Nein'
        score = getattr(project, element)
        comments = getattr(project, f'{element}_comments', '') or ' '
        
        data.append([element_name, questions_answered, score, comments])
    
    table = Table(data, colWidths=[100, 150, 120, 120])
    
    style = TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
        ('ALIGN', (0, 1), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 1), (-1, -1), 10),
        ('TOPPADDING', (0, 1), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ])
    
    table.setStyle(style)
    
    return table

# PDF Export for a project
@routes_bp.route('/export_pdf/<int:project_id>')
def export_pdf(project_id):
    project = Project.query.get_or_404(project_id)
    buffer = BytesIO()
    #doc = SimpleDocTemplate(buffer, pagesize=letter)
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    story = []

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='Justify', alignment=TA_JUSTIFY))

    # Title
    story.append(Paragraph(f"Project: {project.name}", styles['Title']))
    story.append(Spacer(1, 12))

    # Add Scorecard Table
    story.append(Paragraph("MEDDPICC Scorecard", styles['Heading2']))
    story.append(Spacer(1, 12))
    scorecard_table = create_scorecard_table(project)
    story.append(scorecard_table)
    story.append(Spacer(1, 12))

    # Add Radar Chart
    story.append(Paragraph("Radar Chart", styles['Heading2']))
    story.append(Spacer(1, 12))
    radar_path = os.path.join('static', 'images', f'radar_{project_id}.png')
    if os.path.exists(radar_path):
        img = Image(radar_path, width=300, height=300)
        story.append(img)
    else:
        story.append(Paragraph("Radar chart image not found.", styles['Normal']))
    story.append(Spacer(1, 12))

    # Add Question Catalog
    story.append(Paragraph("Question Catalog", styles['Heading2']))
    story.append(Spacer(1, 12))

    elements = ['metrics', 'economic_buyer', 'decision_criteria', 'decision_process', 'paper_process',
                'implications_of_pain', 'champion', 'competition']

    for element in elements:
        story.append(Paragraph(element.replace('_', ' ').capitalize(), styles['Heading3']))
        story.append(Spacer(1, 6))
        
        for q_num in [1, 2, 3]:
            question_field = f"{element}_question{q_num}"
            question_text = getattr(project, question_field)
            story.append(Paragraph(f"Q{q_num}: {question_text}", styles['Normal']))
            story.append(Spacer(1, 3))
        
        story.append(Paragraph(f"Score: {getattr(project, element)}", styles['Normal']))
        story.append(Spacer(1, 12))

    doc.build(story)
    buffer.seek(0)
    return send_file(buffer, as_attachment=True, download_name=f'project_{project_id}.pdf', mimetype='application/pdf')

# Neue Funktion für PowerPoint-Export
def add_scorecard_pptx(slide, project):
    elements = ['metrics', 'economic_buyer', 'decision_criteria', 'decision_process', 'paper_process',
                'implications_of_pain', 'champion', 'competition']
    
    rows, cols = len(elements) + 1, 4  # +1 for header row
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.4 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(3.5)
    
    # Add header row
    header_cells = table.rows[0].cells
    header_texts = ['Element', 'Fragen beantwortet', 'Score (1-10)', 'Kommentare']
    for i, text in enumerate(header_texts):
        cell = header_cells[i]
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add data rows
    for row, element in enumerate(elements, start=1):
        element_name = element.replace('_', ' ').capitalize()
        questions_answered = 'Ja' if getattr(project, f'{element}_question_answered', False) else 'Nein'
        score = getattr(project, element)
        comments = getattr(project, f'{element}_comments', '') or ' '
        
        cells = table.rows[row].cells
        cells[0].text = element_name
        cells[1].text = questions_answered
        cells[2].text = str(score)
        cells[3].text = comments
        
        # Format cells
        for cell in cells:
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.alignment = PP_ALIGN.LEFT
        
        # Color score cell based on value
        score_cell = cells[2]
        if score >= 7:
            score_cell.fill.solid()
            score_cell.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green
        elif score <= 4:
            score_cell.fill.solid()
            score_cell.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
        else:
            score_cell.fill.solid()
            score_cell.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange
    
    return table

# Export Powerpoint pptx
@routes_bp.route('/export_ppt/<int:project_id>')
def export_ppt(project_id):
    project = Project.query.get_or_404(project_id)
    
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Project: {project.name}"
    subtitle.text = "MEDDPICC Analysis"

    # Scorecard Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Manually add title to the slide
    left = top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_shape.text_frame.text = "MEDDPICC Scorecard"
    title_shape.text_frame.paragraphs[0].font.size = Pt(24)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_scorecard_pptx(slide, project)

    # Radar Chart Slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Manually add title to the slide
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_shape.text_frame.text = "Radar Chart"
    title_shape.text_frame.paragraphs[0].font.size = Pt(24)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    radar_path = os.path.join('static', 'images', f'radar_{project_id}.png')
    if os.path.exists(radar_path):
        left = Inches(1.5)
        top = Inches(1.5)
        width = height = Inches(6)
        slide.shapes.add_picture(radar_path, left, top, width, height)
    else:
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Radar chart image not found."

    # Question Catalog Slides
    elements = ['metrics', 'economic_buyer', 'decision_criteria', 'decision_process', 'paper_process',
                'implications_of_pain', 'champion', 'competition']

    for element in elements:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Manually add title to the slide
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.text = element.replace('_', ' ').capitalize()
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for q_num in [1, 2, 3]:
            question_field = f"{element}_question{q_num}"
            question_text = getattr(project, question_field)
            p = tf.add_paragraph()
            p.text = f"Q{q_num}: {question_text}"
            p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"Score: {getattr(project, element)}"
        p.font.bold = True

    # Save the presentation
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return send_file(ppt_buffer, as_attachment=True, 
                     download_name=f'project_{project_id}.pptx', 
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
